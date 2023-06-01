import os
import comtypes.client
import sendmail


def ppt_to_pdf(inp, out, email):
    input_path = os.getcwd() + fr'\ppts\{inp}'
    output_path = os.getcwd() + fr'\pdfs\{out}'
    # try:
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    slides = powerpoint.Presentations.Open(input_path)
    slides.SaveAs(output_path, 32)  # 32 represents the PDF file format
    slides.Close()
    powerpoint.Quit()
    print("pdf generate aytu")
    sendmail.send_mail(email, fr'pdfs\{out}')
    # except Exception as e:
    #     print(e)


from pptx import Presentation


def ppt_to_image():
    ppt_file = "ppts/example.pptx"
    presentation = Presentation(ppt_file)
    output_dir = "imgs"

    for i, slide in enumerate(presentation.slides):
        slide_image = f"{output_dir}/slide_{i + 1}.png"
        slide.export(slide_image)

import pyautogui

def ppt_to_image_():
    ppt_file = "ppts/example.pptx"
    presentation = Presentation(ppt_file)
    output_dir = "imgs"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    for i, slide in enumerate(presentation.slides):
        slide_file = f"{output_dir}/slide_{i + 1}.pptx"
        slide_presentation = Presentation()
        slide_presentation.slides.add_slide(slide)
        slide_presentation.save(slide_file)

        # Convert the slide file to an image using pyautogui
        slide_image = f"{output_dir}/slide_{i + 1}.png"
        os.startfile(slide_file)
        pyautogui.sleep(2)  # Wait for the slide to open
        pyautogui.hotkey('alt', 'f')
        pyautogui.press('d')
        pyautogui.typewrite(slide_image)
        pyautogui.press('enter')
        pyautogui.sleep(2)  # Wait for the image to be saved
        pyautogui.hotkey('alt', 'f4')  # Close the slide file

        # Delete the temporary slide file
        os.remove(slide_file)

ppt_to_image_()
