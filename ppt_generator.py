import os

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
import sendmail
from docx2pdf import convert


def replace_ppt_text(replacements, email):
    ppt_file = r'ppts/example.pptx'
    prs = Presentation(ppt_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for old_word, new_word in replacements:
                            if old_word in run.text:
                                run.text = run.text.replace(old_word, new_word)
                                paragraph.alignment = PP_ALIGN.CENTER
    prs.save(f'ppts\\{replacements[0][1]}.pptx')
    convert(f'ppts\\{replacements[0][1]}.pptx', f'ppts\\{replacements[0][1]}.pdf')
    print('ppt generated successfully ')
    # sendmail.send_mail(email, fr'ppts\{replacements[0][1]}.pptx')
    # ppt_to_pdf.ppt_to_pdf(f'{replacements[0][1]}.pptx', f'{replacements[0][1]}.pdf', email)


replacements = [('{name}', 'row[0]'), ('{college}', 'row[1')]
replace_ppt_text(replacements, "")
